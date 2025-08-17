import io
import json
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import uuid

# Document parsing libraries
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

from app.models import DocumentType

logger = logging.getLogger(__name__)


class DocumentChunk:
    def __init__(self, content: str, metadata: Dict[str, Any]):
        self.content = content
        self.metadata = metadata


class DocumentParser:
    """Base class for document parsing"""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at word boundary
            if end < len(text):
                # Find the last space before the end
                last_space = text.rfind(' ', start, end)
                if last_space > start:
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - chunk_overlap
            if start >= len(text):
                break
                
        return chunks


class PDFParser(DocumentParser):
    """PDF document parser with OCR support"""
    
    @staticmethod
    def extract_text(file_path: str) -> List[DocumentChunk]:
        """Extract text from PDF with page-by-page processing"""
        chunks = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        # Extract text from page
                        text = page.extract_text()
                        
                        # If no text found, try OCR (for scanned PDFs)
                        if not text.strip():
                            logger.info(f"No text found on page {page_num}, attempting OCR")
                            text = PDFParser._ocr_page(page)
                        
                        if text.strip():
                            # Split page text into chunks
                            page_chunks = PDFParser.chunk_text(text)
                            
                            for chunk_idx, chunk_text in enumerate(page_chunks):
                                chunks.append(DocumentChunk(
                                    content=chunk_text,
                                    metadata={
                                        "page_number": page_num,
                                        "chunk_index": chunk_idx,
                                        "source_type": "text_extraction"
                                    }
                                ))
                    except Exception as e:
                        logger.error(f"Error processing page {page_num}: {e}")
                        continue
                        
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise
            
        return chunks
    
    @staticmethod
    def _ocr_page(page) -> str:
        """Perform OCR on a PDF page"""
        try:
            # This is a simplified OCR implementation
            # In production, you might want to use pdf2image + pytesseract
            logger.warning("OCR functionality needs pdf2image library for full implementation")
            return ""
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""


class WordParser(DocumentParser):
    """Microsoft Word document parser"""
    
    @staticmethod
    def extract_text(file_path: str) -> List[DocumentChunk]:
        """Extract text from Word document"""
        chunks = []
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract text from paragraphs
            full_text = ""
            paragraphs_with_styles = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text += para.text + "\n"
                    paragraphs_with_styles.append({
                        "text": para.text,
                        "style": para.style.name if para.style else "Normal"
                    })
            
            # Extract text from tables
            table_text = ""
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text += " | ".join(row_text) + "\n"
            
            if table_text:
                full_text += "\nTables:\n" + table_text
            
            # Create chunks
            text_chunks = WordParser.chunk_text(full_text)
            
            for chunk_idx, chunk_text in enumerate(text_chunks):
                chunks.append(DocumentChunk(
                    content=chunk_text,
                    metadata={
                        "chunk_index": chunk_idx,
                        "source_type": "word_document",
                        "has_tables": len(doc.tables) > 0
                    }
                ))
                
        except Exception as e:
            logger.error(f"Error reading Word document {file_path}: {e}")
            raise
            
        return chunks


class ExcelParser(DocumentParser):
    """Excel spreadsheet parser"""
    
    @staticmethod
    def extract_text(file_path: str) -> List[DocumentChunk]:
        """Extract text from Excel spreadsheet"""
        chunks = []
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Process sheet without pandas
                sheet_text = f"Sheet: {sheet_name}\n\n"
                
                # Get all rows
                all_rows = list(sheet.iter_rows(values_only=True))
                
                if not all_rows:
                    continue
                
                # Filter out completely empty rows
                non_empty_rows = []
                for row in all_rows:
                    if any(cell is not None and str(cell).strip() for cell in row):
                        non_empty_rows.append(row)
                
                if not non_empty_rows:
                    continue
                
                # Add first row as headers if it contains text
                first_row = non_empty_rows[0]
                headers = [str(cell) if cell is not None else "" for cell in first_row]
                if any(header.strip() for header in headers):
                    sheet_text += "Headers: " + " | ".join(h for h in headers if h.strip()) + "\n\n"
                    data_rows = non_empty_rows[1:]
                else:
                    data_rows = non_empty_rows
                
                # Add data rows
                for row in data_rows:
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text += " | ".join(row_text) + "\n"
                
                # Create chunks for this sheet
                if sheet_text.strip():
                    text_chunks = ExcelParser.chunk_text(sheet_text)
                    
                    for chunk_idx, chunk_text in enumerate(text_chunks):
                        chunks.append(DocumentChunk(
                            content=chunk_text,
                            metadata={
                                "sheet_name": sheet_name,
                                "chunk_index": chunk_idx,
                                "source_type": "excel_sheet",
                                "total_sheets": len(workbook.sheetnames)
                            }
                        ))
                        
        except Exception as e:
            logger.error(f"Error reading Excel file {file_path}: {e}")
            raise
            
        return chunks


class PowerPointParser(DocumentParser):
    """PowerPoint presentation parser"""
    
    @staticmethod
    def extract_text(file_path: str) -> List[DocumentChunk]:
        """Extract text from PowerPoint presentation"""
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"Slide {slide_num}:\n\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                # Extract text from notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text += f"\nNotes: {notes_text}\n"
                
                # Create chunks for this slide
                if slide_text.strip() and slide_text != f"Slide {slide_num}:\n\n":
                    text_chunks = PowerPointParser.chunk_text(slide_text)
                    
                    for chunk_idx, chunk_text in enumerate(text_chunks):
                        chunks.append(DocumentChunk(
                            content=chunk_text,
                            metadata={
                                "slide_number": slide_num,
                                "chunk_index": chunk_idx,
                                "source_type": "powerpoint_slide",
                                "total_slides": len(prs.slides)
                            }
                        ))
                        
        except Exception as e:
            logger.error(f"Error reading PowerPoint file {file_path}: {e}")
            raise
            
        return chunks


class DocumentParserService:
    """Main service for parsing different document types"""
    
    PARSERS = {
        DocumentType.PDF: PDFParser,
        DocumentType.DOCX: WordParser,
        DocumentType.XLSX: ExcelParser,
        DocumentType.PPTX: PowerPointParser,
    }
    
    @classmethod
    def parse_document(cls, file_path: str, file_type: DocumentType) -> List[DocumentChunk]:
        """Parse document based on file type"""
        parser = cls.PARSERS.get(file_type)
        
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        logger.info(f"Parsing {file_type} document: {file_path}")
        
        try:
            chunks = parser.extract_text(file_path)
            logger.info(f"Successfully parsed document into {len(chunks)} chunks")
            return chunks
        except Exception as e:
            logger.error(f"Failed to parse document {file_path}: {e}")
            raise
    
    @classmethod
    def get_supported_types(cls) -> List[str]:
        """Get list of supported document types"""
        return [doc_type.value for doc_type in cls.PARSERS.keys()]